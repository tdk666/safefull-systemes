import os
try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

data_dir = "."
output_file = "extracted_content.txt"

with open(output_file, "w") as out:
    for filename in sorted(os.listdir(data_dir)):
        filepath = os.path.join(data_dir, filename)
        if filename.endswith(".pdf") or filename.endswith(".pptx") or filename.endswith(".ppt"):
            out.write(f"\n\n--- {filename} ---\n")
            try:
                if filename.endswith(".pdf") and PdfReader:
                    reader = PdfReader(filepath)
                    for page in reader.pages:
                        text = page.extract_text()
                        if text:
                            out.write(text + "\n")
                elif filename.endswith(".pptx") and Presentation:
                    prs = Presentation(filepath)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                out.write(shape.text + "\n")
                else:
                    out.write("Skipped or format not supported by script directly (like .ppt)\n")
            except Exception as e:
                out.write(f"Error: {e}\n")
print("Done extracting.")
